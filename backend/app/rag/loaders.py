"""
Document loaders: turn an uploaded file on disk into plain text.
Each loader returns (pages, was_ocr_used) where pages is a list of
(text, page_number) tuples, so downstream chunking keeps page-level
citation metadata and the caller can flag documents that needed OCR.
"""
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from app.core.config import settings
from app.core.exceptions import ValidationAppError
from app.core.logging_config import logger
from app.rag import ocr as ocr_module


def load_pdf(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    """
    Extracts text per page with pdfplumber. Any page with too little text
    (heuristically: a scanned/image-only page) falls back to OCR via
    PyMuPDF + EasyOCR, unless OCR is disabled in settings.
    """
    pages: list[tuple[str, int | None]] = []
    used_ocr = False

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = (page.extract_text() or "").strip()

            if settings.OCR_ENABLED and ocr_module.is_scanned_page(text):
                try:
                    ocr_text = ocr_module.ocr_pdf_page(path, i - 1)
                    if len(ocr_text.strip()) > len(text):
                        text = ocr_text.strip()
                        used_ocr = True
                except Exception as exc:  # pragma: no cover
                    logger.warning(f"OCR failed for page {i} of {path}: {exc}")

            if text:
                pages.append((text, i))

    return pages, used_ocr


def load_docx(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    doc = DocxDocument(path)
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [(full_text, None)], False


def load_text(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    text = Path(path).read_text(encoding="utf-8", errors="ignore")
    return [(text, None)], False


def load_csv(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    import csv

    lines = []
    with open(path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        header = next(reader, [])
        for row in reader:
            pairs = ", ".join(f"{h}: {v}" for h, v in zip(header, row))
            lines.append(pairs)
    return [("\n".join(lines), None)], False


def load_excel(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    workbook = load_workbook(path, data_only=True, read_only=True)
    pages: list[tuple[str, int | None]] = []
    for sheet_idx, sheet in enumerate(workbook.worksheets, start=1):
        rows_text = []
        header = None
        for row in sheet.iter_rows(values_only=True):
            values = ["" if v is None else str(v) for v in row]
            if header is None:
                header = values
                continue
            pairs = ", ".join(f"{h}: {v}" for h, v in zip(header, values) if v)
            if pairs:
                rows_text.append(pairs)
        if rows_text:
            pages.append((f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text), sheet_idx))
    return pages, False


def load_pptx(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    prs = Presentation(path)
    pages: list[tuple[str, int | None]] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(t for t in texts if t.strip())
        if slide_text.strip():
            pages.append((slide_text, slide_idx))
    return pages, False


def load_image(path: str) -> tuple[list[tuple[str, int | None]], bool]:
    if not settings.OCR_ENABLED:
        raise ValidationAppError("Image uploads require OCR, which is currently disabled.")
    text = ocr_module.ocr_image_file(path)
    if not text.strip():
        raise ValidationAppError("OCR could not extract any text from this image.")
    return [(text, None)], True


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".txt": load_text,
    ".md": load_text,
    ".csv": load_csv,
    ".xlsx": load_excel,
    ".pptx": load_pptx,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
}


def load_document(path: str, file_type: str) -> tuple[list[tuple[str, int | None]], bool]:
    loader = LOADER_MAP.get(file_type.lower())
    if not loader:
        raise ValidationAppError(f"Unsupported file type: {file_type}")
    pages, used_ocr = loader(path)
    if not pages or all(not text.strip() for text, _ in pages):
        raise ValidationAppError(
            "No extractable text found in document. It may be empty, corrupted, or (if OCR is "
            "disabled) a scanned file."
        )
    return pages, used_ocr
